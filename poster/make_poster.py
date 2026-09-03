#!/usr/bin/env python3
"""
Midwest Population Genetics Meeting 2026 (UChicago) -- poster template for project#3.

Builds poster/MWPG2026_poster_template.pptx : a 30 x 40 inch portrait canvas in the
three-zone architecture (full-width header / two-column core / full-width resolution
banner), and prepares downsampled PNG assets from the 300-dpi TIFFs in main/.

FRAMING AUTHORITY is CLAUDE.md "Project framing" + Amendment F3, NOT scratch.tex's
Introduction (Rule 10). Specifically:
  * The hero claim is the settled A-claim, and it owns the two-branch disjunction
    (not-heritable / not-directional) instead of attributing the excess to balancing
    or fluctuating selection -- the bound is indifferent between the branches and
    holds either way.
  * Lewontin and Buffalo do not appear as thesis or foil, only in the reference strip.
  * The unlinked-only result leads; the linked/chromosomal machinery is marked
    SUPPORTING.
Numbers that scratch.tex still carries as [XX] are rendered as [FILL: ...] markers.
Nothing numeric is invented.

LAYOUT is measured, not guessed: text heights come from the real font metrics via
PIL, blocks are stacked by a flow cursor, and figures contain-fit whatever space is
left over. That makes text/figure collision structurally impossible. A validation
pass at the end reports any region that still came up short.

Run from the repo root:  python poster/make_poster.py
"""

from __future__ import annotations

import glob
import os
import re
import sys

from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

Image.MAX_IMAGE_PIXELS = None

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(HERE, "assets")
FIGDIR = os.path.join(ROOT, "main")
POSTERFIGS = os.path.join(ROOT, "poster", "figs")
OUT = os.path.join(HERE, "MWPG2026_poster_template.pptx")

WARNINGS: list[str] = []
PLACED: list[tuple] = []   # (label, placed_w, placed_h, native_w, native_h, scale)

# ---------------------------------------------------------------- palette ----
PARCHMENT = RGBColor(0xF9, 0xF9, 0xFB)   # canvas background
SLATE     = RGBColor(0x1E, 0x24, 0x2B)   # primary text / structure
INDIGO    = RGBColor(0x2C, 0x4A, 0x6F)   # subheadings, neutral baseline
RUST      = RGBColor(0xC8, 0x52, 0x27)   # accent / discrepancy signal
SAGE      = RGBColor(0x4D, 0x7A, 0x68)   # secondary data / mediators
CARD      = RGBColor(0xFF, 0xFF, 0xFF)   # card / panel fill
BORDER    = RGBColor(0xE5, 0xE7, 0xEB)   # card border
MUTED     = RGBColor(0x6B, 0x72, 0x80)   # captions, de-emphasis
TINT_I    = RGBColor(0xF4, 0xF6, 0xF9)   # indigo wash (equation bands)
TINT_S    = RGBColor(0xF3, 0xF7, 0xF5)   # sage wash

# ------------------------------------------------------------- typography ----
F_DISPLAY = "Inter"              # title + section headers
F_BODY    = "Source Sans 3"      # body copy + captions
F_MATH    = "Latin Modern Math"  # display equations
F_MONO    = "Fira Code"          # data annotations

# The user's hierarchy, scaled x1.6 for a true 24x36 print (it was x2 at 30x40;
# the sheet lost 20% of its width, so the type comes down with it). Body at 26 pt
# and captions at 18 pt are both inside the normal poster range.
S_TITLE   = 76
S_AUTHOR  = 37
S_AFFIL   = 24
S_SECTION = 34    # panel headline
S_SUBHEAD = 32    # letter chip
S_BODY    = 26
S_CALLOUT = 32    # verdicts / hero data callouts
S_HERO    = 34    # header claim, primary line
S_SMALL   = 20    # legends, citations
S_CAP     = 18    # figure captions
S_EYEBROW = 22
LINE_H    = 1.35
LINE_H_TIGHT = 1.22

# ------------------------------------------------------------ page layout ----
PW, PH = 24.0, 36.0
MARGIN, TOP = 1.00, 0.90
CW = PW - 2 * MARGIN             # 22.00
GUTTER = 0.60
COLW = (CW - GUTTER) / 2         # 10.70
LX = MARGIN
RX = MARGIN + COLW + GUTTER

ZONE_GAP = 0.40
Z3_BOTTOM = 34.00                # resolution banner ends here
FOOT_Y, FOOT_H = 34.45, 0.65     # -> 35.10, leaving a 0.90 bottom margin
PAD = 0.32                       # inner card padding
MIN_FIG_H = 2.00                 # warn if a figure region gets thinner than this


# ========================================================= text measurement ==
_FONT_FILES = {
    F_DISPLAY: ["Inter-VariableFont*.ttf"],
    F_BODY:    ["SourceSans3-VariableFont*.ttf"],
    F_MATH:    ["latinmodern-math.otf"],
    F_MONO:    ["FiraCode-VariableFont*.ttf"],
}
_FONT_DIRS = [
    os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts"),
    os.path.join(os.environ.get("LOCALAPPDATA", ""), "Microsoft", "Windows", "Fonts"),
]
_REF_PX = 200.0            # measure once at this size, scale linearly
_BOLD_FUDGE = 1.045        # bold advance vs regular, when variations unavailable
_WRAP_SAFETY = 0.955       # measure against a slightly narrower box than PowerPoint
_cache: dict[tuple, ImageFont.FreeTypeFont] = {}


def _font(name: str, bold: bool):
    key = (name, bold)
    if key in _cache:
        return _cache[key]
    path = None
    for pat in _FONT_FILES.get(name, []):
        for d in _FONT_DIRS:
            if d:
                hit = glob.glob(os.path.join(d, pat))
                if hit:
                    path = hit[0]
                    break
        if path:
            break
    f = None
    if path:
        try:
            f = ImageFont.truetype(path, int(_REF_PX))
            if bold:
                try:
                    f.set_variation_by_name("Bold")
                except Exception:
                    pass
        except Exception:
            f = None
    _cache[key] = f
    return f


_MARKUP = re.compile(r"\*\*|_\{|\^\{|\}")


def plain(s: str) -> str:
    """Strip the mini-markup so the string can be measured."""
    return _MARKUP.sub("", s)


def _adv(s: str, name: str, size_pt: float, bold: bool) -> float:
    """Advance width of `s` in points."""
    f = _font(name, bold)
    if f is None:
        return len(s) * size_pt * 0.50
    w = f.getlength(s) * size_pt / _REF_PX
    return w * (_BOLD_FUDGE if bold and not _variable_ok(name) else 1.0)


def _variable_ok(name: str) -> bool:
    f = _font(name, True)
    return bool(f and hasattr(f, "set_variation_by_name"))


def _words(s: str, bold: bool):
    """Split marked-up text into (word, is_bold) plus "\\n" break tokens.

    Bold spans must be measured with the bold face or a `**...**` run silently
    gains a line at render time, which is what collides with the next block.
    """
    out, cur, b, i, n = [], "", bold, 0, len(s)

    def flush():
        nonlocal cur
        for wd in cur.split():
            out.append((wd, b))
        cur = ""

    while i < n:
        if s.startswith("**", i):
            flush()
            b = not b
            i += 2
        elif s.startswith("_{", i) or s.startswith("^{", i):
            j = s.index("}", i)
            cur += s[i + 2:j]                      # sub/sup renders smaller: safe
            i = j + 1
        elif s[i] == "\n":
            flush()
            out.append(("\n", b))
            i += 1
        else:
            cur += s[i]
            i += 1
    flush()
    return out


def wrap(s: str, name: str, size_pt: float, width_in: float, bold: bool = False):
    """Greedy, bold-aware word wrap. Returns the laid-out lines."""
    limit = width_in * 72.0 * _WRAP_SAFETY
    space = _adv(" ", name, size_pt, bold)
    lines, cur, cur_w = [], "", 0.0
    for wd, wb in _words(s, bold):
        if wd == "\n":
            lines.append(cur)
            cur, cur_w = "", 0.0
            continue
        ww = _adv(wd, name, size_pt, wb)
        add = ww if not cur else ww + space
        if cur and cur_w + add > limit:
            lines.append(cur)
            cur, cur_w = wd, ww
        else:
            cur = wd if not cur else cur + " " + wd
            cur_w += add
    lines.append(cur)
    return lines


def h_of(s: str, name: str, size_pt: float, width_in: float, *,
         line_h: float = LINE_H, bold: bool = False) -> float:
    """Rendered height in inches."""
    n = len(wrap(s, name, size_pt, width_in, bold))
    return n * size_pt * line_h / 72.0


def fit_size(s: str, name: str, width_in: float, height_in: float, start_pt: float,
             *, line_h: float = LINE_H, bold: bool = False, floor_pt: float = 14.0):
    """Largest size <= start_pt whose wrapped height fits height_in."""
    size = start_pt
    while size > floor_pt and h_of(s, name, size, width_in, line_h=line_h,
                                  bold=bold) > height_in:
        size -= 0.5
    return size


# ============================================================== primitives ===
def rect(shapes, x, y, w, h, fill=None, line=None, line_w=1.5,
         shape=MSO_SHAPE.RECTANGLE, dash=None, radius=None):
    s = shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
        if dash is not None:
            s.line.dash_style = dash
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    s.text_frame.text = ""
    return s


def card(shapes, x, y, w, h):
    """Cool-alabaster panel: white fill, hairline border, no shadow."""
    return rect(shapes, x, y, w, h, fill=CARD, line=BORDER, line_w=1.75)


def _emit(p, s, *, font, size, color, bold=False, italic=False):
    """Mini-markup:  **bold**   _{subscript}   ^{superscript}   (nesting-safe)."""
    i, n, b, base = 0, len(s), bold, None
    buf = ""

    def flush():
        nonlocal buf
        if not buf:
            return
        r = p.add_run()
        r.text = buf
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.italic = italic
        r.font.color.rgb = color
        if base is not None:
            # run.font._element *is* the <a:rPr>; set the attribute on it directly.
            r.font._element.set("baseline", base)
        buf = ""

    while i < n:
        if s.startswith("**", i):
            flush()
            b = not b
            i += 2
        elif s.startswith("_{", i) or s.startswith("^{", i):
            flush()
            base = "-25000" if s[i] == "_" else "30000"
            j = s.index("}", i)
            buf = s[i + 2:j]
            flush()
            base = None
            i = j + 1
        else:
            buf += s[i]
            i += 1
    flush()


def text(shapes, x, y, w, h, paras, *, font=F_BODY, size=S_BODY, color=SLATE,
         bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_h=LINE_H, space_after=0):
    """`paras` is a str, or a list of str / (str, overrides)."""
    tb = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(paras, str):
        paras = [paras]
    for i, item in enumerate(paras):
        body, ov = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        # EXACT spacing in points, not a float multiple. A float multiple is applied
        # to the font's own line height (~1.2 x em), which renders ~20% taller than
        # h_of() predicts and silently collides with the next block.
        p.line_spacing = Pt(ov.get("size", size) * ov.get("line_h", line_h))
        p.space_after = Pt(ov.get("space_after", space_after))
        p.space_before = Pt(ov.get("space_before", 0))
        _emit(p, body, font=ov.get("font", font), size=ov.get("size", size),
              color=ov.get("color", color), bold=ov.get("bold", bold),
              italic=ov.get("italic", italic))
    return tb


class Flow:
    """Vertical cursor: every block reports its measured height and advances y."""

    def __init__(self, shapes, x, y, w):
        self.sh, self.x, self.y, self.w = shapes, x, y, w

    def gap(self, dy):
        self.y += dy
        return self

    def para(self, paras, *, size=S_BODY, font=F_BODY, color=SLATE, line_h=LINE_H,
             bold=False, align=PP_ALIGN.LEFT, lead=0.20, w=None):
        w = self.w if w is None else w
        items = [paras] if isinstance(paras, str) else paras
        h = 0.0
        for it in items:
            body, ov = it if isinstance(it, tuple) else (it, {})
            h += h_of(body, ov.get("font", font), ov.get("size", size), w,
                      line_h=ov.get("line_h", line_h), bold=ov.get("bold", bold))
        h += lead * (len(items) - 1)
        text(self.sh, self.x, self.y, w, h + 0.06, items, font=font, size=size,
             color=color, bold=bold, line_h=line_h, align=align,
             space_after=lead * 72)
        self.y += h + 0.06
        return self

    def h_para(self, paras, *, size=S_BODY, font=F_BODY, line_h=LINE_H, bold=False,
               lead=0.20, w=None):
        w = self.w if w is None else w
        items = [paras] if isinstance(paras, str) else paras
        h = sum(h_of(it[0] if isinstance(it, tuple) else it,
                     (it[1] if isinstance(it, tuple) else {}).get("font", font),
                     (it[1] if isinstance(it, tuple) else {}).get("size", size), w,
                     line_h=(it[1] if isinstance(it, tuple) else {}).get("line_h", line_h),
                     bold=(it[1] if isinstance(it, tuple) else {}).get("bold", bold))
                for it in items)
        return h + lead * (len(items) - 1)

    def para_fit(self, paras, avail_h, *, start=S_BODY, floor=18.0, label="text", **kw):
        """Largest body size <= start whose flowed height fits avail_h."""
        size = start
        while size > floor and self.h_para(paras, size=size, **kw) > avail_h:
            size -= 0.5
        if self.h_para(paras, size=size, **kw) > avail_h + 0.02:
            WARNINGS.append("%s: does not fit even at %.0f pt" % (label, size))
        elif size < start - 0.01:
            print("  autoshrink %-22s %.0f -> %.0f pt" % (label, start, size))
        self.para(paras, size=size, **kw)
        return size

    def rule(self, thickness=0.035, color=BORDER, inset=0.0):
        rect(self.sh, self.x + inset, self.y, self.w - 2 * inset, thickness, fill=color)
        self.y += thickness
        return self


def eyebrow(shapes, x, y, w, label, color=INDIGO, size=S_EYEBROW):
    h = size * 1.30 / 72.0
    tb = text(shapes, x, y, w, h, label.upper(), font=F_DISPLAY, size=size,
              color=color, bold=True, line_h=1.30)
    tb.text_frame.paragraphs[0].runs[0].font._element.set("spc", "320")
    return h


def panel_header(flow, letter, title, accent=INDIGO, note=None):
    """Accent rule + letter chip + headline (+ optional right-aligned label)."""
    sh, x, y, w = flow.sh, flow.x, flow.y, flow.w
    rect(sh, x, y, w, 0.055, fill=accent)
    chip = 0.68
    ty = y + 0.28
    nw = (_adv(note, F_DISPLAY, S_SMALL, True) / 72.0 + 0.62) if note else 0.0
    tw = w - chip - 0.28 - (nw + 0.28 if note else 0.0)
    th = max(chip, h_of(title, F_DISPLAY, S_SECTION, tw, line_h=1.08, bold=True))
    rect(sh, x, ty + (th - chip) / 2, chip, chip, fill=accent,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.16)
    text(sh, x, ty + (th - chip) / 2, chip, chip, letter, font=F_DISPLAY,
         size=S_SUBHEAD, color=CARD, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
    text(sh, x + chip + 0.28, ty, tw, th, title, font=F_DISPLAY, size=S_SECTION,
         color=SLATE, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_h=1.08)
    if note:
        nh = S_SMALL * 1.30 / 72.0 + 0.24
        rect(sh, x + w - nw, ty + (th - nh) / 2, nw, nh, fill=TINT_I,
             line=RGBColor(0xDD, 0xE3, 0xEB), line_w=1.25)
        text(sh, x + w - nw, ty + (th - nh) / 2, nw, nh, note, font=F_DISPLAY,
             size=S_SMALL, color=INDIGO, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    flow.y = ty + th + 0.30
    return flow


def verdict(flow, body, *, color=RUST, size=S_CALLOUT, w=None):
    """Accent rule + bold callout. Advances the flow cursor."""
    w = flow.w if w is None else w
    bar, gap = 0.14, 0.24
    tw = w - bar - gap
    h = h_of(body, F_DISPLAY, size, tw, line_h=LINE_H_TIGHT, bold=True) + 0.10
    rect(flow.sh, flow.x, flow.y, bar, h, fill=color)
    text(flow.sh, flow.x + bar + gap, flow.y, tw, h, body, font=F_DISPLAY, size=size,
         color=color, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_h=LINE_H_TIGHT)
    flow.y += h
    return flow


def equation(flow, body, *, size=38, w=None):
    w = flow.w if w is None else w
    h = h_of(body, F_MATH, size, w - 0.5, line_h=1.20) + 0.44
    rect(flow.sh, flow.x, flow.y, w, h, fill=TINT_I,
         line=RGBColor(0xDD, 0xE3, 0xEB), line_w=1.25)
    text(flow.sh, flow.x + 0.25, flow.y, w - 0.5, h, body, font=F_MATH, size=size,
         color=INDIGO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_h=1.20)
    flow.y += h
    return flow


def badge(flow, body, *, color=SAGE, fill=TINT_S, size=S_SMALL, w=None):
    w = flow.w if w is None else w
    h = h_of(body, F_BODY, size, w - 0.52, line_h=LINE_H_TIGHT) + 0.34
    rect(flow.sh, flow.x, flow.y, w, h, fill=fill, line=color, line_w=1.25)
    text(flow.sh, flow.x + 0.26, flow.y, w - 0.52, h, body, size=size, color=color,
         anchor=MSO_ANCHOR.MIDDLE, line_h=LINE_H_TIGHT)
    flow.y += h
    return flow


def tag(flow, label, *, color=INDIGO):
    size = S_SMALL
    w = _adv(plain(label), F_DISPLAY, size, True) / 72.0 + 0.62
    h = size * 1.30 / 72.0 + 0.26
    rect(flow.sh, flow.x, flow.y, w, h, fill=TINT_I,
         line=RGBColor(0xDD, 0xE3, 0xEB), line_w=1.25)
    text(flow.sh, flow.x, flow.y, w, h, label, font=F_DISPLAY, size=size, color=color,
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    flow.y += h
    return flow


def dashed_note(flow, paras, *, color=SAGE, size=None, w=None, avail_h=None,
                label="note"):
    w = flow.w if w is None else w
    size = size or S_BODY * 0.85
    iw = w - 0.68

    def _h(sz):
        return (sum(h_of(p, F_BODY, sz, iw, line_h=1.30) for p in paras)
                + 0.18 * (len(paras) - 1) + 0.52)

    if avail_h is not None:
        while size > 17.0 and _h(size) > avail_h:
            size -= 0.5
        if _h(size) > avail_h + 0.02:
            WARNINGS.append("%s: does not fit even at %.0f pt" % (label, size))
    h = _h(size) if avail_h is None else max(_h(size), avail_h)
    rect(flow.sh, flow.x, flow.y, w, h, fill=RGBColor(0xF5, 0xF8, 0xF6), line=color,
         line_w=2.0, dash=MSO_LINE_DASH_STYLE.DASH)
    text(flow.sh, flow.x + 0.34, flow.y + 0.26, iw, h - 0.52, list(paras), size=size,
         color=color, line_h=1.30, space_after=13)
    flow.y += h
    return flow


def caption(flow, body, *, w=None):
    w = flow.w if w is None else w
    h = h_of(body, F_BODY, S_CAP, w, line_h=1.28)
    text(flow.sh, flow.x, flow.y, w, h + 0.06, body, size=S_CAP, color=MUTED,
         line_h=1.28)
    flow.y += h + 0.06
    return flow


# --------------------------------------------------------------- artwork -----
def prepare_assets(specs):
    """Resolve each slot to a TIFF, preferring poster/figs, and convert to PNG."""
    os.makedirs(ASSETS, exist_ok=True)
    out = {}
    for slot, (candidates, max_w) in specs.items():
        src = origin = None
        for name in candidates:
            for d, tag in ((POSTERFIGS, "poster-scale"), (FIGDIR, "MANUSCRIPT-SCALE")):
                cand = os.path.join(d, name + ".tiff")
                if os.path.exists(cand):
                    src, origin, base = cand, tag, name
                    break
            if src:
                break
        if src is None:
            WARNINGS.append("slot %s: no figure found (%s) -- renders as a FILL box"
                            % (slot, ", ".join(candidates)))
            continue
        im = Image.open(src).convert("RGB")
        target = min(im.width, int(round(max_w * 300)))
        if target < im.width:
            im = im.resize((target, int(round(im.height * target / im.width))),
                           Image.LANCZOS)
        dst = os.path.join(ASSETS, "panel_%s.png" % slot)
        im.save(dst, "PNG", optimize=True)
        out[slot] = dst
        flag = "" if origin == "poster-scale" else "   <-- fallback"
        print("  %s: %-42s %-16s %4dx%-4d px%s"
              % (slot, base[:42], origin, im.width, im.height, flag))
        if origin != "poster-scale":
            WARNINGS.append("slot %s uses the journal-page figure; run "
                            "scripts/R/poster_figs.R for the poster-scale one" % slot)
    return out


def figure_block(flow, path, avail_h, cap, *, label="figure", w=None):
    """Contain-fit an image into the space left over, then caption it."""
    w = flow.w if w is None else w
    cap_h = h_of(cap, F_BODY, S_CAP, w, line_h=1.28) + 0.20 if cap else 0.0
    box_h = avail_h - cap_h
    if box_h < MIN_FIG_H:
        WARNINGS.append("%s: only %.2f in of figure height (min %.2f)"
                        % (label, box_h, MIN_FIG_H))
    if path is None:
        rect(flow.sh, flow.x, flow.y, w, max(box_h, 0.6), fill=TINT_S, line=SAGE,
             line_w=2.0, dash=MSO_LINE_DASH_STYLE.DASH)
        text(flow.sh, flow.x, flow.y, w, max(box_h, 0.6), "[FILL: %s]" % label,
             font=F_MONO, size=S_SMALL, color=SAGE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        used_h = max(box_h, 0.6)
    else:
        iw, ih = Image.open(path).size
        scale = min(w / iw, box_h / ih)
        fw, fh = iw * scale, ih * scale
        flow.sh.add_picture(path, Inches(flow.x + (w - fw) / 2), Inches(flow.y),
                            Inches(fw), Inches(fh))
        used_h = fh
        # Report placed vs native size: a scale well below 1.0 means the
        # figure's own type is being shrunk, which is what kills legibility.
        PLACED.append((label, fw, fh, iw / 300.0, ih / 300.0, fw / (iw / 300.0)))
    flow.y += used_h + (0.20 if cap else 0.0)
    if cap:
        caption(flow, cap, w=w)
    return flow


# ================================================================= content ===
TITLE = "Measured Fitness Variance Is Too Weak to Explain the Missing Diversity"
AUTHOR = "Walid Mawass"
AFFIL = ("[FILL: Department / Committee]  ·  University of Chicago  ·  "
         "mawass@uchicago.edu")

HERO_CLAIM = (
    "Heritable variance in fitness erodes neutral diversity, and a pedigree measures that "
    "variance without touching a single base. The 19 published estimates remove at most "
    "twentyfold of it. The deficit these species carry runs far deeper.")
# Short form in the header; the mechanism is unpacked in Zone 3 next to the two
# boxes, which is where it belongs. F3 requires the caveat wherever the bound is
# stated -- this is that caveat, compressed to one line.
HERO_BRANCHES = ("Diversity also sets an **upper limit** on sustained fitness variance. Every "
                 "value reported so far stays under it.")

# Logical panel slot -> (candidate basenames in priority order, max placed width).
# poster/figs holds versions rendered at their exact placed size by
# scripts/R/poster_figs.R, so their axis type is correct on the wall; main/ holds
# the journal-page versions, kept only as a fallback.
FIGSPEC = {
    "A": (["PosterFig_A_Baseline_Constraint", "MainFig_Baseline_Constraint"], 6.0),
    "B": (["PosterFig_B_Bound_vs_Delta"], 10.1),
    "C": (["PosterFig_C_Mu_Sensitivity",
           "MainFig_Empirical_Sensitivity_MutationRate"], 10.1),
    "D": (["PosterFig_D_Va_Deflation", "SuppFig3_Va_Deflation_Analysis"], 6.0),
}


# =================================================================== build ===
def build():
    print("preparing assets from main/ ...")
    figs = prepare_assets(FIGSPEC)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(PW), Inches(PH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sh = slide.shapes
    rect(sh, 0, 0, PW, PH, fill=PARCHMENT)

    # Zone 1 height and Zone 3 height are both derived from measured content;
    # Zone 2 takes whatever is between them.
    z1_h = zone1(sh, TOP)
    z3_h = zone3(sh, measure_only=True)
    z3_y = Z3_BOTTOM - z3_h
    z2_y = TOP + z1_h + ZONE_GAP
    z2_h = (z3_y - ZONE_GAP) - z2_y
    ph = (z2_h - ZONE_GAP) / 2
    print("\nzones: Z1 %.2f in | Z2 %.2f in (panels %.2f) | Z3 %.2f in"
          % (z1_h, z2_h, ph, z3_h))
    if ph < 6.8:
        WARNINGS.append("core panels only %.2f in tall" % ph)

    panel_a(sh, LX, z2_y, COLW, ph, figs)
    panel_b(sh, LX, z2_y + ph + ZONE_GAP, COLW, ph, figs)
    panel_c(sh, RX, z2_y, COLW, ph, figs)
    panel_d(sh, RX, z2_y + ph + ZONE_GAP, COLW, ph, figs)
    zone3(sh, y=z3_y, h=z3_h)
    footer(sh)

    prs.save(OUT)
    print("\nwrote %s  (%.0f x %.0f in)" % (OUT, PW, PH))
    if PLACED:
        print("\nFIGURE PLACEMENT  (scale = placed in / native in at 300 dpi)")
        for lab, fw, fh, nw, nh, sc in PLACED:
            flag = "  <-- type shrunk" if sc < 0.92 else ""
            print("  %-42s placed %5.2f x %-5.2f  native %5.2f x %-5.2f  scale %.2f%s"
                  % (lab[:42], fw, fh, nw, nh, sc, flag))
    if WARNINGS:
        print("\nLAYOUT WARNINGS")
        for w in WARNINGS:
            print("  -", w)
    else:
        print("layout check: clean -- every block fits its region")


# ------------------------------------------------- Zone 1: header + claim ----
def zone1(sh, y):
    f = Flow(sh, MARGIN, y, CW)
    f.para(TITLE, font=F_DISPLAY, size=S_TITLE, bold=True, line_h=1.03,
           align=PP_ALIGN.CENTER)
    f.gap(0.18)
    f.para([(AUTHOR, dict(font=F_DISPLAY, size=S_AUTHOR, bold=True)),
            (AFFIL, dict(size=S_AFFIL, color=INDIGO))],
           align=PP_ALIGN.CENTER, lead=0.10, line_h=1.22)
    f.gap(0.26)

    # Hero card: measured, then drawn behind its own contents.
    hx, hy = MARGIN, f.y
    iw = CW - 0.22 - 2 * PAD
    eb = S_EYEBROW * 1.30 / 72.0
    h1 = h_of(HERO_CLAIM, F_DISPLAY, S_HERO, iw, line_h=1.20, bold=True)
    h2 = h_of(HERO_BRANCHES, F_BODY, S_BODY, iw, line_h=1.30)
    hh = 0.30 + eb + 0.14 + h1 + 0.20 + h2 + 0.30

    card(sh, hx, hy, CW, hh)
    rect(sh, hx, hy, 0.22, hh, fill=RUST)
    g = Flow(sh, hx + 0.22 + PAD, hy + 0.30, iw)
    g.gap(eyebrow(sh, g.x, g.y, iw, "The claim") + 0.14)
    g.para(HERO_CLAIM, font=F_DISPLAY, size=S_HERO, color=RUST, bold=True, line_h=1.20)
    g.gap(0.20)
    g.para(HERO_BRANCHES, size=S_BODY, line_h=1.30)
    return (hy + hh) - y


# --------------------------------------- Panel A: the diversity-free bound ---
A_BODY = [
    "Because a neutral site's association with its background halves each generation "
    "rather than vanishing, heritable differences in offspring number erode diversity "
    "under free recombination. **No linkage map, no annotation, no equilibrium "
    "assumption.**",
]
A_CAP = ("**Fig. A** · Diversity kept, against no selection at the same census size "
         "(log). Dashed: 1. Dotted: 0.1. **No input is calibrated on diversity.**")


def panel_a(sh, x, y, w, h, figs):
    card(sh, x, y, w, h)
    f = Flow(sh, x + PAD, y + PAD, w - 2 * PAD)
    panel_header(f, "A", "Diversity erodes without linkage")
    top, avail = f.y, (y + h - PAD) - f.y

    # The figure is portrait (0.80) and therefore height-bound: it claims the full
    # column height, and the text column takes the width left over.
    fig_h = avail - h_of(A_CAP, F_BODY, S_CAP, 5.0, line_h=1.28) - 0.26
    fw = min(4.70, fig_h * 0.80 + 0.02)
    tw = f.w - fw - 0.34

    vtxt, vsize = "At most 3×; 1.4× at the median", 27
    t = Flow(sh, x + PAD, top, tw)
    equation(t, "N_{e}/N  =  1 / (1 + 4κV_{A})", size=34)
    t.gap(0.24)
    reserve = (0.28 + h_of(vtxt, F_DISPLAY, vsize, tw - 0.38, line_h=LINE_H_TIGHT,
                           bold=True) + 0.10)
    t.para_fit(A_BODY, (y + h - PAD) - t.y - reserve, lead=0.20, label="panel A body")
    t.gap(0.28)
    verdict(t, vtxt, size=vsize)

    g = Flow(sh, x + PAD + tw + 0.34, top, fw)
    figure_block(g, figs.get("A"), avail, A_CAP, label="Panel A figure")
    _check("panel A text", t.y, y + h - PAD)
    _check("panel A figure", g.y, y + h - PAD)


# ------------------------------------- Panel B: the per-species test (TBD) ---
B_BODY = [
    "Inverting the same expression gives that limit: **V_{A,max} = (1 − δ)/(4κδ)**, "
    "where **δ is the fraction of neutral diversity still standing**. Any estimate of "
    "fitness variance has to come in under it.",
]
B_CAP = ("**Fig. B** · Curves: the largest V_{A} each mating system allows, with "
         "N the number of breeding adults. Rust: reported pedigree V_{A}, quartiles "
         "and one tick per population (n = 19). [FILL: observed δ]")
B_VERDICT = "At δ = 0.1 the limit is 2.25; nothing reaches it."


def panel_b(sh, x, y, w, h, figs):
    card(sh, x, y, w, h)
    f = Flow(sh, x + PAD, y + PAD, w - 2 * PAD)
    panel_header(f, "B", "Diversity sets the upper limit", accent=SAGE)
    bottom = y + h - PAD

    vh = h_of(B_VERDICT, F_DISPLAY, S_CALLOUT * 0.86, f.w - 0.38,
              line_h=LINE_H_TIGHT, bold=True) + 0.10

    f.para_fit(B_BODY, 1.45, lead=0.20, label="panel B body")
    f.gap(0.24)
    figure_block(f, figs.get("B"), bottom - f.y - (0.26 + vh), B_CAP,
                 label="Panel B figure")
    f.gap(0.26)
    verdict(f, B_VERDICT, size=S_CALLOUT * 0.86)
    _check("panel B", f.y, bottom)


# ------------------------------------ Panel C: linkage only tightens it ------
C_BODY = [
    "Near a selected site that association survives far longer than one generation of "
    "recombination allows. Adding recombination maps and functional target sizes at the "
    "chromosome level only tightens the estimate: **a further 1.2 to 10×**.",
]
C_CAP = ("**Fig. C** · (left) Diversity kept with the linked term, across the 95% "
         "interval on μ. (right) Genome-wide linkage penalty Ω̄. Needs maps and "
         "annotations; Panel A needs neither.")
C_VERDICT = "Total: twentyfold, at the outside."


def panel_c(sh, x, y, w, h, figs):
    card(sh, x, y, w, h)
    f = Flow(sh, x + PAD, y + PAD, w - 2 * PAD)
    panel_header(f, "C", "Linkage only tightens it", note="MAP-DEPENDENT")
    bottom = y + h - PAD

    f.para_fit(C_BODY, 1.95, lead=0.18, label="panel C body")
    f.gap(0.24)
    figure_block(f, figs.get("C"), bottom - f.y, C_CAP, label="Panel C figure")
    _check("panel C", f.y, bottom)


# --------------------------------------------- Panel D: the partition -------
D_BODY = [
    "Even counting **all** of the reported V_{A} as sustained, the prediction stays "
    "above the diversity these populations actually keep. Counting less of it widens the "
    "gap rather than closing it.",
    "Whatever is depleting their diversity, **this is not it**.",
]
D_CAP = ("**Fig. D** · Diversity kept when 100 / 50 / 25 / 10% of reported V_{A} is "
         "counted as sustained. Dotted: [FILL: observed δ].")
D_VERDICT = "The gap widens.\nIt does not close."


def panel_d(sh, x, y, w, h, figs):
    card(sh, x, y, w, h)
    f = Flow(sh, x + PAD, y + PAD, w - 2 * PAD)
    panel_header(f, "D", "Far short of the missing diversity", accent=RUST)
    top, avail = f.y, (y + h - PAD) - f.y

    fig_h = avail - h_of(D_CAP, F_BODY, S_CAP, 5.0, line_h=1.28) - 0.26
    fw = min(4.70, fig_h * 0.80 + 0.02)
    tw = f.w - fw - 0.34

    t = Flow(sh, x + PAD, top, tw)
    reserve = 0.30 + h_of(D_VERDICT, F_DISPLAY, S_CALLOUT, tw - 0.38,
                          line_h=LINE_H_TIGHT, bold=True) + 0.10
    t.para_fit(D_BODY, avail - reserve, lead=0.20, label="panel D body")
    t.gap(0.30)
    verdict(t, D_VERDICT)

    g = Flow(sh, x + PAD + tw + 0.34, top, fw)
    figure_block(g, figs.get("D"), avail, D_CAP, label="Panel D figure")
    _check("panel D text", t.y, y + h - PAD)
    _check("panel D figure", g.y, y + h - PAD)


# ------------------------------- Zone 3: mechanism + takeaway for the field --
Z3_HEAD = "What these models need V_{A} to be"
Z3_LEDE = ("The calculation above and Fisher's theorem draw on one quantity: **fitness "
           "variance that is transmitted and sustained.**")
Z3_BOXES = [
    (INDIGO, "WHAT THE EQUATIONS NEED", "sustained over many generations",
     "Neutral diversity responds only to fitness variance **maintained** over coalescent "
     "time. A pedigree gives you a few decades.",
     "▸ A pedigree cannot show it."),
    (SAGE, "A SECOND ROUTE, NO DIVERSITY NEEDED", "mutation and selection alone",
     "Mutation and selection can sustain only so much variance. That route asks the same "
     "question with no diversity data, once the scales are reconciled.",
     "▸ Independent of everything left."),
]
Z3_TAKEAWAY = [
    "Any estimate of fitness variance, from a pedigree or from sequence, has to come in "
    "under (1 − δ)/(4κδ). **That limit is the reusable part of this.**",
    "Reported values remove at most twentyfold of neutral diversity, and nothing on the "
    "input side comes from diversity data. **The number is the finding, not the "
    "inequality.**",
    "Large measured V_{A} gets read as large adaptive potential. Yet that reading needs "
    "the variance to persist, and nothing here shows it does — **a discussion point, "
    "not a claim.**",
]
Z3_VERDICT = "Ceiling at observed δ: [FILL: V_{A,max}]."


def zone3(sh, y=None, h=None, measure_only=False):
    w = CW
    lw = w * 0.545 - PAD
    rw = w - 2 * PAD - lw - 0.55
    bw = (lw - 0.34) / 2
    tsize = S_BODY * 0.88

    head_h = h_of(Z3_HEAD, F_DISPLAY, S_SECTION * 1.05, w - 2 * PAD, line_h=1.10,
                  bold=True)
    lede_h = h_of(Z3_LEDE, F_BODY, S_BODY * 0.92, w - 2 * PAD, line_h=1.30)
    box_body = max(h_of(b[3], F_BODY, S_BODY * 0.84, bw - 0.60, line_h=1.28)
                   for b in Z3_BOXES)
    box_tag = max(h_of(b[4], F_DISPLAY, S_SMALL, bw - 0.60, line_h=1.24, bold=True)
                  for b in Z3_BOXES)
    box_note = max(h_of(b[2], F_BODY, S_SMALL * 0.94, bw - 0.60, line_h=1.24, bold=True)
                   for b in Z3_BOXES)
    box_h = (0.58 + 0.24 + 0.20 + 0.20 + box_note + 0.26 + box_body + 0.18
             + box_tag + 0.22)
    vh = h_of(Z3_VERDICT, F_DISPLAY, S_CALLOUT * 0.78, lw - 0.38,
              line_h=LINE_H_TIGHT, bold=True) + 0.10
    left_h = box_h + 0.24 + vh
    right_h = (S_EYEBROW * 1.30 / 72.0 + 0.18
               + sum(h_of(pp, F_BODY, tsize, rw, line_h=1.30) for pp in Z3_TAKEAWAY)
               + 0.20 * (len(Z3_TAKEAWAY) - 1))
    need = PAD + head_h + 0.12 + lede_h + 0.26 + max(left_h, right_h) + PAD + 0.14

    if measure_only:
        return need

    x = MARGIN
    card(sh, x, y, w, h)
    rect(sh, x, y, w, 0.055, fill=SLATE)
    f = Flow(sh, x + PAD, y + PAD, w - 2 * PAD)
    f.para(Z3_HEAD, font=F_DISPLAY, size=S_SECTION * 1.05, bold=True, line_h=1.10)
    f.gap(0.12)
    f.para(Z3_LEDE, size=S_BODY * 0.92, line_h=1.30)
    f.gap(0.26)
    cy = f.y

    for i, (accent, head, note, body, tagline) in enumerate(Z3_BOXES):
        bx = x + PAD + i * (bw + 0.34)
        rect(sh, bx, cy, bw, box_h, fill=RGBColor(0xFC, 0xFD, 0xFD), line=accent,
             line_w=2.0)
        rect(sh, bx, cy, bw, 0.58, fill=accent)
        text(sh, bx + 0.18, cy, bw - 0.36, 0.58, head, font=F_DISPLAY, size=S_SMALL,
             color=CARD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        b = Flow(sh, bx + 0.30, cy + 0.58 + 0.24, bw - 0.60)

        # schematic: a chromosome block with the selected sites picked out
        rect(sh, b.x + 0.04, b.y, b.w - 0.08, 0.20, fill=RGBColor(0xD5, 0xD9, 0xDE),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
        n = 7
        step = (b.w - 0.08 - 0.30) / (n - 1)
        for k in range(n):
            sel = (k % 3 == 1)
            d = 0.30 if sel else 0.18
            rect(sh, b.x + 0.04 + 0.15 + k * step - d / 2, b.y + 0.10 - d / 2, d, d,
                 fill=(accent if sel else RGBColor(0x9C, 0xA3, 0xAF)),
                 shape=MSO_SHAPE.OVAL)
        b.gap(0.20 + 0.20)
        b.para(note, size=S_SMALL * 0.94, color=accent, bold=True,
               align=PP_ALIGN.CENTER, line_h=1.24)
        b.gap(0.26 - 0.06)
        b.para(body, size=S_BODY * 0.84, line_h=1.28)
        b.gap(0.18)
        b.para(tagline, font=F_DISPLAY, size=S_SMALL, color=accent, bold=True,
               line_h=1.24)
        _check("zone 3 box %d" % (i + 1), b.y, cy + box_h)

    v = Flow(sh, x + PAD, cy + box_h + 0.24, lw)
    verdict(v, Z3_VERDICT, size=S_CALLOUT * 0.78)

    r = Flow(sh, x + PAD + lw + 0.55, cy, rw)
    r.gap(eyebrow(sh, r.x, r.y, rw, "Takeaway for the field") + 0.18)
    r.para_fit(Z3_TAKEAWAY, (y + h - PAD) - r.y, start=tsize, line_h=1.30, lead=0.20,
               label="zone 3 takeaway")
    _check("zone 3 takeaway", r.y, y + h - PAD)
    _check("zone 3 verdict", v.y, y + h - PAD)
    return need


def footer(sh):
    rect(sh, MARGIN, FOOT_Y - 0.18, CW, 0.035, fill=BORDER)
    text(sh, MARGIN, FOOT_Y, CW * 0.68, FOOT_H,
         "Robertson 1961 · Santiago & Caballero 1995 · Bonnet et al. 2022 · "
         "Corbett-Detig et al. 2015 · Buffalo 2021 · Lewin & Eyre-Walker 2026 · "
         "Connallon & Czuppon 2026 · Charlesworth 2026 · Eyre-Walker & Keightley 2007 · "
         "Huber et al. 2017",
         size=S_SMALL * 0.90, color=MUTED, line_h=1.22)
    text(sh, MARGIN + CW * 0.70, FOOT_Y, CW * 0.30, FOOT_H,
         "Midwest Population Genetics Meeting 2026 · University of Chicago\n"
         "[FILL: QR / preprint link]",
         size=S_SMALL * 0.90, color=MUTED, align=PP_ALIGN.RIGHT, line_h=1.22)


def _check(label, got_bottom, limit_bottom):
    if got_bottom > limit_bottom + 0.10:
        WARNINGS.append("%s overflows by %.2f in" % (label, got_bottom - limit_bottom))


if __name__ == "__main__":
    sys.exit(build())
